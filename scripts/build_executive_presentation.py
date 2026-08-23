"""
ADPilot Pro v3.0 - Professional Executive Presentation Generator (Enhanced Visuals & Diagrams)
Embeds high-resolution AI visuals, system diagrams, UI preview mockups,
and mathematical formulations onto 28 custom-styled widescreen slides.
"""

import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette Constants ──
BG_DARK = RGBColor(7, 9, 14)         # #07090E (Obsidian Dark)
BG_CARD = RGBColor(13, 19, 34)       # #0D1322 (Elevated Card Dark)
BG_CARD_SEC = RGBColor(10, 15, 29)   # #0A0F1D (Secondary Card)
BG_CARD_ACCENT = RGBColor(16, 24, 45)# #10182D (Border Highlight)

TEXT_WHITE = RGBColor(248, 250, 252) # #F8FAFC (Pure White)
TEXT_MUTED = RGBColor(148, 163, 184) # #94A3B8 (Cool Gray)
TEXT_DIM = RGBColor(100, 116, 139)   # #64748B (Slate Dim)

CYAN_PRIMARY = RGBColor(0, 240, 255) # #00F0FF (Electric Cyan)
CYAN_DARK = RGBColor(6, 182, 212)    # #06B6D4 (Teal)
EMERALD_GREEN = RGBColor(16, 185, 129)# #10B981 (Matrix Emerald)
PURPLE_ACCENT = RGBColor(139, 92, 246)# #8B5CF6 (Vivid Purple)
AMBER_WARNING = RGBColor(245, 158, 11)# #F59E0B (Amber Flame)
ROSE_DANGER = RGBColor(244, 63, 94)  # #F43F5E (Rose Red)
BLUE_SUPPORT = RGBColor(59, 130, 246) # #3B82F6 (Vivid Blue)

FONT_HEADING = "Segoe UI"
FONT_BODY = "Segoe UI"
FONT_MONO = "Consolas"

ASSETS_DIR = r"Presentation\assets"

def get_asset(filename):
    p = os.path.join(ASSETS_DIR, filename)
    return p if os.path.exists(p) else None

def create_presentation():
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_header_footer(slide, category_text, slide_title, subtitle_text, slide_num, total_slides=28):
        # 1. Background Canvas
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()

        # Cyber Grid Background
        for i in range(1, 14):
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(i), 0, Inches(0.01), Inches(7.5))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(15, 23, 42)
            line.line.fill.background()
        for i in range(1, 8):
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(i), Inches(13.333), Inches(0.01))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(15, 23, 42)
            line.line.fill.background()

        # 2. Header Section
        # Category Badge
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.0), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = FONT_MONO
        p_cat.font.size = Pt(9.5)
        p_cat.font.bold = True
        p_cat.font.color.rgb = CYAN_PRIMARY

        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.68), Inches(9.5), Inches(0.55))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = slide_title
        p_title.font.name = FONT_HEADING
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE

        # Subtitle Context
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.22), Inches(10.5), Inches(0.35))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle_text
        p_sub.font.name = FONT_BODY
        p_sub.font.size = Pt(10.5)
        p_sub.font.color.rgb = TEXT_MUTED

        # Slide Index Badge (Top Right)
        idx_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.2), Inches(0.48), Inches(1.333), Inches(0.4))
        idx_card.fill.solid()
        idx_card.fill.fore_color.rgb = BG_CARD
        idx_card.line.color.rgb = RGBColor(30, 41, 59)
        tf_idx = idx_card.text_frame
        p_idx = tf_idx.paragraphs[0]
        p_idx.alignment = PP_ALIGN.CENTER
        p_idx.text = f"{slide_num:02d} / {total_slides}"
        p_idx.font.name = FONT_MONO
        p_idx.font.size = Pt(10)
        p_idx.font.bold = True
        p_idx.font.color.rgb = CYAN_PRIMARY

        # 3. Footer Bar
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.95), Inches(11.733), Inches(0.015))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(30, 41, 59)
        line.line.fill.background()

        foot_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.02), Inches(11.733), Inches(0.3))
        tf_foot = foot_box.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = "ADPILOT PRO v3.0  •  AUTONOMOUS AI MARKETING OPERATING SYSTEM  •  ACADEMIC CAPSTONE DEFENSE 2026"
        p_foot.font.name = FONT_MONO
        p_foot.font.size = Pt(8)
        p_foot.font.color.rgb = TEXT_DIM

    def draw_card(slide, left, top, width, height, title="", border_color=None, bg_color=None):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color or BG_CARD
        card.line.color.rgb = RGBColor(30, 41, 59)
        card.line.width = Pt(1.0)

        if border_color:
            # HUD Cyber Accents (Top & Left)
            accent_top = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.06))
            accent_top.fill.solid()
            accent_top.fill.fore_color.rgb = border_color
            accent_top.line.fill.background()
            
            accent_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(0.06), Inches(height))
            accent_left.fill.solid()
            accent_left.fill.fore_color.rgb = border_color
            accent_left.line.fill.background()

        if title:
            title_box = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + 0.2), Inches(width - 0.5), Inches(0.35))
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = title.upper()
            p.font.name = FONT_MONO
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = border_color or CYAN_PRIMARY
        return card

    def embed_framed_image(slide, left, top, width, height, image_path, caption=None, border_color=CYAN_PRIMARY):
        if image_path and os.path.exists(image_path):
            # Frame card
            frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
            frame.fill.solid()
            frame.fill.fore_color.rgb = BG_CARD
            frame.line.color.rgb = RGBColor(30, 41, 59)
            frame.line.width = Pt(1.0)

            # HUD Cyber Accents (Top & Left)
            accent_top = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.06))
            accent_top.fill.solid()
            accent_top.fill.fore_color.rgb = border_color
            accent_top.line.fill.background()
            
            accent_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(0.06), Inches(height))
            accent_left.fill.solid()
            accent_left.fill.fore_color.rgb = border_color
            accent_left.line.fill.background()

            # Insert Image
            img_margin = 0.15
            slide.shapes.add_picture(
                image_path,
                Inches(left + img_margin),
                Inches(top + img_margin),
                Inches(width - 2 * img_margin),
                Inches(height - (0.45 if caption else 2 * img_margin))
            )

            # Optional caption bar
            if caption:
                cbox = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + height - 0.4), Inches(width - 0.2), Inches(0.35))
                tf = cbox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.text = caption.upper()
                p.font.name = FONT_MONO
                p.font.size = Pt(9)
                p.font.bold = True
                p.font.color.rgb = border_color

    # =========================================================================
    # SLIDE 1: COVER SLIDE (With AI Hero Graphic)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = BG_DARK
    bg1.line.fill.background()

    top_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.04))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = CYAN_PRIMARY
    top_bar.line.fill.background()

    # Left Column: Brand & Details (6.2 inches)
    pill1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.9), Inches(4.5), Inches(0.38))
    pill1.fill.solid()
    pill1.fill.fore_color.rgb = BG_CARD
    pill1.line.color.rgb = CYAN_PRIMARY
    pill1.text_frame.paragraphs[0].text = "ENTERPRISE AI OPERATING SYSTEM  //  v3.0"
    pill1.text_frame.paragraphs[0].font.name = FONT_MONO
    pill1.text_frame.paragraphs[0].font.size = Pt(9.5)
    pill1.text_frame.paragraphs[0].font.bold = True
    pill1.text_frame.paragraphs[0].font.color.rgb = CYAN_PRIMARY

    h1_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.45), Inches(5.8), Inches(1.1))
    tf1 = h1_box.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "ADPILOT PRO"
    p1.font.name = FONT_HEADING
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE

    sub1_box = s1.shapes.add_textbox(Inches(0.8), Inches(2.55), Inches(5.8), Inches(0.9))
    tf1_sub = sub1_box.text_frame
    tf1_sub.word_wrap = True
    p1_sub = tf1_sub.paragraphs[0]
    p1_sub.text = "Autonomous Intelligence for Strategic Planning, Creative Synthesis, PPO Reinforcement Learning & Deterministic Decision Governance"
    p1_sub.font.name = FONT_BODY
    p1_sub.font.size = Pt(12.5)
    p1_sub.font.color.rgb = TEXT_MUTED

    # 4 Highlights Grid on Left
    left_feats = [
        ("18-STAGE DAG", "Pydantic v2 typed contracts", CYAN_PRIMARY),
        ("PPO RL OPTIMIZER", "12-dim continuous budget policy", EMERALD_GREEN),
        ("DUAL-STREAM RAG", "FastEmbed BGE + Qdrant (MRR 1.0)", PURPLE_ACCENT),
        ("HITL AUDIT LEDGER", "HMAC-SHA256 cryptographic sign", AMBER_WARNING)
    ]
    for i, (title, desc, col) in enumerate(left_feats):
        c_idx = i % 2
        r_idx = i // 2
        cl = 0.8 + c_idx * 2.9
        ct = 3.6 + r_idx * 1.15
        draw_card(s1, cl, ct, 2.75, 1.0, title, border_color=col)
        tbox = s1.shapes.add_textbox(Inches(cl + 0.3), Inches(ct + 0.5), Inches(2.3), Inches(0.55))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = FONT_BODY
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MUTED

    # Right Column: Hero Visual Image (5.5 inches wide)
    embed_framed_image(
        s1,
        left=6.85,
        top=0.9,
        width=5.68,
        height=4.85,
        image_path=get_asset("hero_cover.jpg") or get_asset("hero_banner.png"),
        caption="Autonomous Multi-Agent Neural Command Architecture",
        border_color=CYAN_PRIMARY
    )

    # Footer Metadata
    meta_box = s1.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(11.733), Inches(0.8))
    tf_meta = meta_box.text_frame
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = "PROJECT DEFENSE: Digital Pioneers Initiative (DiGiLiANS)  •  MCIT Egypt  •  Military Technical College (MTC)  •  Capstone 2026"
    p_meta.font.name = FONT_MONO
    p_meta.font.size = Pt(9.5)
    p_meta.font.bold = True
    p_meta.font.color.rgb = CYAN_PRIMARY

    p_meta2 = tf_meta.add_paragraph()
    p_meta2.text = "CORE ENGINEERING: Mohamed Gharieb (Lead), Ahmed Awni, Mohamed Sleem, Khaled Mohamed, Karem Tarek"
    p_meta2.font.name = FONT_BODY
    p_meta2.font.size = Pt(9.5)
    p_meta2.font.color.rgb = TEXT_DIM

    # =========================================================================
    # SLIDE 2: EXECUTIVE AGENDA
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    add_header_footer(s2, "STRATEGIC ROADMAP", "Executive Presentation Agenda", "Structured walkthrough of the market crisis, multi-agent AI architecture, mathematical models, empirical proof, and commercial ROI.", 2)
    
    agenda_modules = [
        ("01", "The Problem & Market Drain", "The marketing fragmentation crisis, tool silos, and the real cost of operational lag.", CYAN_PRIMARY),
        ("02", "ADPilot Pro Vision & Architecture", "The unified AI OS, master 18-stage execution DAG, and context integrity.", BLUE_SUPPORT),
        ("03", "Autonomous Multi-Agent Fleet", "Specialized agents for Strategy, Research, Content, Vision, Analytics, and Optimization.", PURPLE_ACCENT),
        ("04", "Intelligence Stack & Math Models", "LLM Reasoning vs Specialized ML, continuous PPO RL formulation, and CLIP-ViT vision.", EMERALD_GREEN),
        ("05", "Dual-Stream RAG & Memory", "FastEmbed BGE embeddings, Qdrant vector store, and 4-tier persistent memory.", CYAN_DARK),
        ("06", "Governance, HITL & Simulation", "Cryptographic HMAC-SHA256 review gates and end-to-end campaign execution trace.", AMBER_WARNING),
        ("07", "Empirical Evaluation & UI Experience", "52/52 Vitest & 269 Pytest verification, production cyber UI, and comparative benchmarks.", ROSE_DANGER),
        ("08", "Commercial Value & Defense Close", "Scalable SaaS monetization, limitations transparency, future roadmap, and conclusion.", CYAN_PRIMARY),
    ]

    for i, (num, title, desc, col) in enumerate(agenda_modules):
        col_idx = i % 2
        row_idx = i // 2
        c_left = 0.8 + col_idx * 5.95
        c_top = 1.75 + row_idx * 1.22
        
        draw_card(s2, c_left, c_top, 5.75, 1.08, border_color=RGBColor(30, 41, 59))
        
        nb = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(c_left + 0.3), Inches(c_top + 0.15), Inches(0.75), Inches(0.75))
        nb.fill.solid()
        nb.fill.fore_color.rgb = BG_DARK
        nb.line.color.rgb = col
        nb.text_frame.paragraphs[0].text = num
        nb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        nb.text_frame.paragraphs[0].font.name = FONT_MONO
        nb.text_frame.paragraphs[0].font.size = Pt(14)
        nb.text_frame.paragraphs[0].font.bold = True
        nb.text_frame.paragraphs[0].font.color.rgb = col

        tbox = s2.shapes.add_textbox(Inches(c_left + 1.05), Inches(c_top + 0.12), Inches(4.55), Inches(0.85))
        tf = tbox.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(11.5)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 3: LEADERSHIP & INSTITUTIONAL BACKING
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    add_header_footer(s3, "GOVERNANCE & LEADERSHIP", "Project Leadership & Institutional Backing", "Developed by specialized AI engineers and supported by national digital transformation and academic leadership.", 3)
    
    # Left Card: Engineering Founders
    draw_card(s3, 0.8, 1.75, 6.8, 4.9, "Core Engineering Team", CYAN_PRIMARY)
    founders = [
        ("Mohamed Abd Elsalam Gharieb", "Team Lead • Master Orchestrator & Architecture", "Architecture, Pydantic v2 DAG, EventBus Telemetry, System Integration"),
        ("Ahmed Awni", "AI Engineer • Research Agent & Market Intelligence", "SerpAPI integration, audience slang extraction, semantic grounding"),
        ("Mohamed Sleem", "AI Engineer • Content Agent & Copywriting", "Multi-format copywriting, Google RSAs, LinkedIn thought leadership, email nurture"),
        ("Khaled Mohamed", "ML Engineer • Predictive Analytics & Scoring", "Ridge revenue forecaster, econometric multi-target regression, evaluation rubrics"),
        ("Karem Tarek", "AI Engineer • Design Agent & Vision Studio", "Nano Banana diffusion prompts, CLIP-ViT visual quality scoring, brand safe-zones")
    ]
    for idx, (name, role, details) in enumerate(founders):
        y_pos = 2.15 + idx * 0.88
        ibox = s3.shapes.add_textbox(Inches(1.0), Inches(y_pos), Inches(6.4), Inches(0.8))
        tf = ibox.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = f"▸ {name}"
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE
        p2 = tf.add_paragraph()
        p2.text = f"   {role}"
        p2.font.name = FONT_MONO
        p2.font.size = Pt(9)
        p2.font.bold = True
        p2.font.color.rgb = CYAN_PRIMARY
        p3 = tf.add_paragraph()
        p3.text = f"   {details}"
        p3.font.name = FONT_BODY
        p3.font.size = Pt(8.5)
        p3.font.color.rgb = TEXT_MUTED

    # Right Card: Institutional Backing
    draw_card(s3, 7.8, 1.75, 4.733, 4.9, "Institutional Supervision", PURPLE_ACCENT)
    inst_box = s3.shapes.add_textbox(Inches(8.0), Inches(2.15), Inches(4.333), Inches(4.3))
    tf_inst = inst_box.text_frame
    tf_inst.word_wrap = True

    sections = [
        ("ACADEMIC TRACK", "Artificial Intelligence & Data Science (AI)\nDigital Pioneers Initiative (DiGiLiANS) Capstone 2026", CYAN_PRIMARY),
        ("SUPERVISION", "Under direct supervision of Academic Leadership and Senior Technical Advisory Committees.", EMERALD_GREEN),
        ("SPONSORING ENTITIES", "• Ministry of Communications and Information Technology (MCIT)\n• Military Technical College (MTC)\n• Egyptian Military Academy", AMBER_WARNING),
        ("DEFENSE RIGOR", "Rigorous evaluation covering formal multi-agent contracts, empirical test suites, continuous RL convergence, and zero-hallucination verification.", TEXT_WHITE)
    ]
    for stitle, sdesc, scol in sections:
        p = tf_inst.add_paragraph()
        p.text = stitle
        p.font.name = FONT_MONO
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = scol
        p2 = tf_inst.add_paragraph()
        p2.text = sdesc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(9)
        p2.font.color.rgb = TEXT_MUTED
        tf_inst.add_paragraph().text = ""

    # =========================================================================
    # SLIDE 4: THE PROBLEM: MARKET REALITY
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    add_header_footer(s4, "MARKET DRAIN", "The Modern Marketing Paradox: High Activity, Broken Flow", "Marketing teams deploy 8+ disconnected SaaS tools, yet campaigns take weeks to launch, brand context decays at every handoff, and real-time optimization is virtually impossible.", 4)
    
    problems = [
        ("TOOL SILOS & FRAGMENTATION", "Teams juggle CRM, SEO tools, Ad managers, Figma, and spreadsheets.", "• Strategic context decays by 40% across tool handoffs (Kietzmann et al.).\n• Manual copy-pasting between platforms introduces human error.\n• Zero unified data model linking creative assets to revenue KPIs.", ROSE_DANGER),
        ("CRIPPLING EXECUTION LAG", "Campaign briefs require 10–14 days of back-and-forth iteration.", "• Copywriters wait on research; designers wait on copy; media buyers wait on assets.\n• Ad creative fatigue sets in before human teams can refresh variants.\n• Fleeting market trends and viral opportunities are missed.", AMBER_WARNING),
        ("MONOLITHIC AI FAILURE", "Generic single-prompt LLMs fail on complex enterprise campaigns.", "• Single prompt chatbots suffer context saturation and loss of nuance.\n• Hallucinations compound across multi-step marketing workflows (Xi et al., 2025).\n• Zero automated validation against brand safety rules or visual margins.", BLUE_SUPPORT)
    ]

    for i, (title, subtitle, bullets, col) in enumerate(problems):
        c_left = 0.8 + i * 3.98
        draw_card(s4, c_left, 1.75, 3.75, 4.9, title, col)
        
        tbox = s4.shapes.add_textbox(Inches(c_left + 0.3), Inches(2.2), Inches(3.25), Inches(4.3))
        tf = tbox.text_frame
        tf.word_wrap = True
        p_sub = tf.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.name = FONT_HEADING
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_WHITE

        tf.add_paragraph().text = ""
        for line in bullets.split("\n"):
            p = tf.add_paragraph()
            p.text = line
            p.font.name = FONT_BODY
            p.font.size = Pt(9.5)
            p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 5: EXECUTIVE IMPACT & FINANCIAL DRAIN
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    add_header_footer(s5, "EXECUTIVE IMPACT", "The Compounding Cost of Manual Marketing Inefficiency", "Leaving marketing workflows uncoordinated creates massive hidden costs in operational drag, ad spend waste, and delayed revenue realization.", 5)
    
    metrics = [
        ("10–14 Days", "AVERAGE CYCLE TIME", "Brief-to-launch delay across copywriting, creative generation, and review loops.", ROSE_DANGER),
        ("68%", "BRAND VOICE DRIFT", "Enterprise marketers reporting inconsistent brand messaging across channels (Dwivedi 2023).", AMBER_WARNING),
        ("20–35%", "BUDGET LEAKAGE", "Ad spend wasted on underperforming ad variants before human teams detect decay.", CYAN_PRIMARY),
        ("70%+", "ADMINISTRATIVE DRAG", "Senior strategist hours consumed by manual formatting, copy-pasting, and tool coordination.", PURPLE_ACCENT)
    ]

    for i, (val, label, desc, col) in enumerate(metrics):
        c_left = 0.8 + i * 2.98
        draw_card(s5, c_left, 1.75, 2.78, 2.6, border_color=col)
        
        vbox = s5.shapes.add_textbox(Inches(c_left + 0.3), Inches(1.95), Inches(2.28), Inches(0.8))
        tf = vbox.text_frame
        p = tf.paragraphs[0]
        p.text = val
        p.font.name = FONT_HEADING
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = col

        lbox = s5.shapes.add_textbox(Inches(c_left + 0.3), Inches(2.75), Inches(2.28), Inches(1.4))
        tf_l = lbox.text_frame
        tf_l.word_wrap = True
        p1 = tf_l.paragraphs[0]
        p1.text = label
        p1.font.name = FONT_MONO
        p1.font.size = Pt(9.5)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE

        p2 = tf_l.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(9)
        p2.font.color.rgb = TEXT_MUTED

    draw_card(s5, 0.8, 4.6, 11.733, 2.05, "The Enterprise Imperative", CYAN_PRIMARY)
    cbox = s5.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(11.333), Inches(1.5))
    tf_c = cbox.text_frame
    tf_c.word_wrap = True
    p = tf_c.paragraphs[0]
    p.text = "Enterprise marketing demands an AI system that does not simply generate text, but acts as a closed-loop operating system:"
    p.font.name = FONT_HEADING
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    reqs = [
        "1. Contextual Integrity: Single structured campaign brief propagating immutably across strategy, copy, visual prompts, and bids.",
        "2. Multi-Specialist Collaboration: Dividing marketing labor among specialized agents rather than relying on one overloaded generalist LLM.",
        "3. Closed-Loop Optimization: Continuously sampling performance signals and dynamically adjusting budget allocations via Reinforcement Learning."
    ]
    for req in reqs:
        p_r = tf_c.add_paragraph()
        p_r.text = f"• {req}"
        p_r.font.name = FONT_BODY
        p_r.font.size = Pt(9.5)
        p_r.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 6: PRODUCT VISION: ONE UNIFIED AI MARKETING OS
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    add_header_footer(s6, "PRODUCT VISION", "The Solution: One Coordinated Marketing Operating System", "ADPilot Pro replaces disconnected toolchains with a single autonomous platform where specialized AI agents orchestrate the entire campaign lifecycle under human governance.", 6)
    
    pillars = [
        ("1. STRUCTURED BRIEF", "Single Source of Truth", "• Ingests business objectives, budget, duration, geography, and tone.\n• Validates input schema via strict Pydantic v2 data models.\n• Eliminates ambiguity before any model execution starts.", CYAN_PRIMARY),
        ("2. MULTI-AGENT DAG", "Specialized Collaboration", "• 18 autonomous agents execute in deterministic dependency order.\n• Asynchronous parallel execution cut runtime by ~40%.\n• Context passed via immutable data contracts with zero drift.", PURPLE_ACCENT),
        ("3. CLOSED-LOOP OPTIMIZER", "Continuous Adaptation", "• Pre-flight computer vision and econometric revenue forecasting.\n• PPO Reinforcement Learning dynamically reallocates budget.\n• Human-in-the-loop governance with cryptographic audit trails.", EMERALD_GREEN)
    ]

    for i, (title, subtitle, bullets, col) in enumerate(pillars):
        c_left = 0.8 + i * 3.98
        draw_card(s6, c_left, 1.75, 3.75, 4.9, title, col)
        
        tbox = s6.shapes.add_textbox(Inches(c_left + 0.3), Inches(2.2), Inches(3.25), Inches(4.3))
        tf = tbox.text_frame
        tf.word_wrap = True
        p_sub = tf.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.name = FONT_HEADING
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_WHITE

        tf.add_paragraph().text = ""
        for line in bullets.split("\n"):
            p = tf.add_paragraph()
            p.text = line
            p.font.name = FONT_BODY
            p.font.size = Pt(9.5)
            p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 7: MASTER PIPELINE: 18-STAGE EXECUTION DAG (With Pipeline Diagram)
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    add_header_footer(s7, "PIPELINE ARCHITECTURE", "Master System Architecture: 18-Stage Execution DAG", "The Master Orchestrator coordinates an 18-stage Directed Acyclic Graph (DAG) executing parallel intelligence, creative synthesis, quality auditing, and RL tuning.", 7)
    
    # Left Column: 4 Execution Phases (5.8 inches)
    stages_group = [
        ("STAGE 01–04: INGESTION & CONTEXT", [("01. Input Ingestion", "Pydantic Brief Validation"), ("02. Product Classifier", "Taxonomy & Vertical ID"), ("03. Tone Extractor", "Brand Voice Calibration"), ("04. Persona Synthesizer", "ICP & Pain Point Profile")], CYAN_PRIMARY),
        ("STAGE 05–08: INTELLIGENCE & RESEARCH", [("05. Strategy Planner", "Cross-Channel Roadmap"), ("06. Audience Research", "FastEmbed BGE Vector Search"), ("07. Competitor Intel", "Battlecards & Market Gaps"), ("08. Content Copywriter", "Multi-Format Copy Variants")], PURPLE_ACCENT),
        ("STAGE 09–13: CREATIVE, CV & OPTIMIZER", [("09. Creative Design", "Nano Banana Visual Prompts"), ("10. CV Quality Gate", "CLIP-ViT Aesthetic Check"), ("11. Predictive Analytics", "Ridge Revenue Forecaster"), ("12. PPO RL Optimizer", "Continuous Action Selection")], EMERALD_GREEN),
        ("STAGE 14–18: GOVERNANCE & EXECUTION", [("13. Constraint Filter", "Budget Simplex Guards"), ("14. HITL Review Gate", "HMAC-SHA256 Sign"), ("15. Publishing Dispatch", "Idempotent Multi-Network API"), ("16–18. Telemetry Stream", "Closed-Loop Feedback")], AMBER_WARNING)
    ]

    for i, (group_title, items, col) in enumerate(stages_group):
        c_idx = i % 2
        r_idx = i // 2
        cl = 0.8 + c_idx * 2.85
        ct = 1.75 + r_idx * 2.45
        draw_card(s7, cl, ct, 2.72, 2.35, group_title, col)
        tbox = s7.shapes.add_textbox(Inches(cl + 0.3), Inches(ct + 0.5), Inches(2.28), Inches(1.9))
        tf = tbox.text_frame
        tf.word_wrap = True
        for s_idx, (st_name, st_desc) in enumerate(items):
            p = tf.add_paragraph() if s_idx > 0 else tf.paragraphs[0]
            p.text = f"▸ {st_name}"
            p.font.name = FONT_HEADING
            p.font.size = Pt(8.5)
            p.font.bold = True
            p.font.color.rgb = TEXT_WHITE
            p2 = tf.add_paragraph()
            p2.text = f"   {st_desc}"
            p2.font.name = FONT_BODY
            p2.font.size = Pt(7.8)
            p2.font.color.rgb = TEXT_MUTED

    # Right Column: High-Resolution DAG Architecture Diagram (5.9 inches)
    embed_framed_image(
        s7,
        left=6.65,
        top=1.75,
        width=5.88,
        height=4.9,
        image_path=get_asset("002.png") or get_asset("pipeline_architecture.png"),
        caption="18-Stage Directed Acyclic Graph (DAG) Execution Topology",
        border_color=CYAN_PRIMARY
    )

    # =========================================================================
    # SLIDE 8: MULTI-AGENT FLEET OVERVIEW (With Agent Fleet Diagram)
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    add_header_footer(s8, "AGENT FLEET", "The Autonomous Multi-Agent Fleet Architecture", "Specialized agents extending the BaseAgent abstract contract pattern, communicating via immutable Pydantic payloads across 4 functional tiers.", 8)
    
    # Left Column: 4 Tiers (5.8 inches)
    agent_tiers = [
        ("TIER 1: STRATEGY & INTEL", [("Strategy Agent", "GPT-4o Router", "Funnel architecture & channel weighting."), ("Research Agent", "FastEmbed BGE", "Live market signals & search intent.")], CYAN_PRIMARY),
        ("TIER 2: CREATIVE & VISION", [("Content Copywriter", "Claude 3.5", "RSAs, LinkedIn posts & email nurture."), ("Design & CV Agent", "Nano Banana + CLIP", "Multi-aspect prompts & aesthetic checks.")], PURPLE_ACCENT),
        ("TIER 3: OPTIMIZATION & RL", [("Analytics Forecaster", "Ridge Regression", "Multi-target ROAS & CAC forecasting."), ("RL Optimizer", "PyTorch PPO", "12-dim continuous budget rebalancing.")], EMERALD_GREEN),
        ("TIER 4: GOVERNANCE & DISPATCH", [("HITL Review Gate", "HMAC-SHA256", "Risk-based human approval ledger."), ("Publishing Agent", "Multi-Network", "Idempotent dispatch & WebSocket telemetry.")], AMBER_WARNING)
    ]

    for i, (tier_title, ag_list, col) in enumerate(agent_tiers):
        c_idx = i % 2
        r_idx = i // 2
        cl = 0.8 + c_idx * 2.85
        ct = 1.75 + r_idx * 2.45
        draw_card(s8, cl, ct, 2.72, 2.35, tier_title, col)
        tbox = s8.shapes.add_textbox(Inches(cl + 0.3), Inches(ct + 0.5), Inches(2.28), Inches(1.9))
        tf = tbox.text_frame
        tf.word_wrap = True
        for a_idx, (ag_name, ag_tech, ag_desc) in enumerate(ag_list):
            p = tf.add_paragraph() if a_idx > 0 else tf.paragraphs[0]
            p.text = f"{ag_name} [{ag_tech}]"
            p.font.name = FONT_HEADING
            p.font.size = Pt(8.5)
            p.font.bold = True
            p.font.color.rgb = col
            p2 = tf.add_paragraph()
            p2.text = ag_desc
            p2.font.name = FONT_BODY
            p2.font.size = Pt(8.0)
            p2.font.color.rgb = TEXT_MUTED

    # Right Column: Agent Fleet Diagram (5.9 inches)
    embed_framed_image(
        s8,
        left=6.65,
        top=1.75,
        width=5.88,
        height=4.9,
        image_path=get_asset("003.png") or get_asset("agent_fleet.png"),
        caption="Autonomous Multi-Agent Collaborative Microservice Topology",
        border_color=PURPLE_ACCENT
    )

    # =========================================================================
    # SLIDE 9: AGENT RESPONSIBILITIES MATRIX
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    add_header_footer(s9, "RESPONSIBILITY MATRIX", "Agent Operational & Technical Contract Matrix", "Strict deterministic interfaces: Every agent ingests a typed Pydantic v2 input contract and emits a verified output guarantee.", 9)
    
    rows = [
        ["AGENT NAME", "CORE RESPONSIBILITY", "TYPED INPUT", "GUARANTEED OUTPUT", "TECH ENGINE"],
        ["Strategy Agent", "Funnel architecture & channel weighting", "CampaignBriefContext", "StrategyPlanContract", "GPT-4o Router"],
        ["Research Agent", "Market trends, search intent & keywords", "TargetAudienceSpec", "MarketIntelContract", "FastEmbed BGE + SerpAPI"],
        ["Competitor Agent", "Rival positioning & counter-hooks", "ProductVerticalSpec", "CompetitorBattlecard", "Qdrant Vector RAG"],
        ["Content Agent", "Multi-channel ad copy & email sequences", "Strategy + Intel", "ContentOutputPackage", "Claude 3.5 Sonnet"],
        ["Design Agent", "Multi-aspect visual synthesis directives", "BrandPalette + Goals", "CreativeAssetPackage", "Nano Banana Studio"],
        ["CV Quality Agent", "Zero-shot visual scoring & margin check", "Generated Image URLs", "CVQualityAuditScore", "CLIP-ViT (ONNX)"],
        ["Analytics Agent", "Econometric revenue & ROAS forecasting", "Strategy + Content", "EconometricForecast", "Scikit-Learn Ridge"],
        ["RL Optimizer", "Dynamic channel budget reallocation", "12-dim State Vector", "ActionVector a_t", "PyTorch PPO Policy"],
        ["HITL Reviewer", "Cryptographic governance & safety sign-off", "High-Risk Actions", "SignedAuditLedger", "HMAC-SHA256"]
    ]

    table_shape = s9.shapes.add_table(10, 5, Inches(0.8), Inches(1.75), Inches(11.733), Inches(4.9))
    table = table_shape.table
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(3.2)
    table.columns[2].width = Inches(2.2)
    table.columns[3].width = Inches(2.4)
    table.columns[4].width = Inches(2.133)

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            if r_idx == 0:
                cell.fill.fore_color.rgb = RGBColor(16, 24, 45)
            else:
                cell.fill.fore_color.rgb = BG_CARD if r_idx % 2 == 0 else BG_CARD_SEC

            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = val
            p.font.name = FONT_MONO if (r_idx == 0 or c_idx in [2, 3, 4]) else FONT_BODY
            p.font.size = Pt(8.5) if r_idx > 0 else Pt(9)
            if r_idx == 0:
                p.font.bold = True
                p.font.color.rgb = CYAN_PRIMARY
            else:
                if c_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = TEXT_WHITE
                elif c_idx == 4:
                    p.font.bold = True
                    p.font.color.rgb = CYAN_PRIMARY
                else:
                    p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 10: LLM REASONING VS SPECIALIZED AI
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    add_header_footer(s10, "CONCEPTUAL FOUNDATION", "LLM Reasoning vs. Specialized Machine Learning", "Why single LLMs fail at enterprise marketing and how ADPilot Pro assigns specialized mathematical roles to appropriate AI architectures.", 10)
    
    comp_models = [
        ("FRONTIER LLMs", "Reasoning & Planning", [("Models", "GPT-4o, Claude 3.5 Sonnet"), ("Role", "Strategic roadmap synthesis, multi-variant copywriting, psychological buying trigger formulation."), ("Why LLM?", "Unmatched language comprehension, high zero-shot reasoning, and creative generation capability.")], CYAN_PRIMARY),
        ("STATISTICAL ML", "Prediction & Scoring", [("Models", "Ridge Regression, TF-IDF Classifier"), ("Role", "Quantitative revenue forecasting, copy quality scoring, historical conversion rate regression."), ("Why ML?", "Deterministic execution in < 3ms, zero hallucination, verifiable statistical confidence intervals.")], EMERALD_GREEN),
        ("REINFORCEMENT LEARNING", "Continuous Optimization", [("Models", "PPO Actor-Critic (PyTorch)"), ("Role", "Sequential channel budget allocation across 12 continuous state variables."), ("Why RL?", "Learns closed-loop decision policies from dynamic reward feedback rather than static rules.")], PURPLE_ACCENT),
        ("COMPUTER VISION", "Perception & Safety", [("Models", "CLIP-ViT (ONNX Runtime)"), ("Role", "Zero-shot visual aesthetics, text clipping prevention, WCAG AAA contrast ratio compliance."), ("Why CV?", "Real-time multimodal inspection verifying creative quality before publishing.")], AMBER_WARNING)
    ]

    for i, (title, subtitle, kvs, col) in enumerate(comp_models):
        c_left = 0.8 + i * 2.98
        draw_card(s10, c_left, 1.75, 2.78, 4.9, title, col)
        tbox = s10.shapes.add_textbox(Inches(c_left + 0.3), Inches(2.2), Inches(2.28), Inches(4.3))
        tf = tbox.text_frame
        tf.word_wrap = True
        p_sub = tf.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.name = FONT_HEADING
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_WHITE
        tf.add_paragraph().text = ""
        for k, v in kvs:
            p_k = tf.add_paragraph()
            p_k.text = k.upper()
            p_k.font.name = FONT_MONO
            p_k.font.size = Pt(8.5)
            p_k.font.bold = True
            p_k.font.color.rgb = col
            p_v = tf.add_paragraph()
            p_v.text = v
            p_v.font.name = FONT_BODY
            p_v.font.size = Pt(8.5)
            p_v.font.color.rgb = TEXT_MUTED
            tf.add_paragraph().text = ""

    # =========================================================================
    # SLIDE 11: CUSTOM ML MODELS REGISTRY
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    add_header_footer(s11, "MODEL REGISTRY", "Custom Trained Machine Learning Model Catalog", "ADPilot Pro maintains an audited artifact registry of custom trained PyTorch weights, Scikit-Learn regressors, and ONNX vision scoring pipelines.", 11)
    
    ml_models = [
        ("PPO POLICY NETWORK", "Reinforcement Learning", "research/models/optimizer/ppo_policy.pt", "PyTorch 2.11", "12-dim Continuous State", "K-dim Action Vector", "15.8ms", "+0.48 Mean Reward", CYAN_PRIMARY),
        ("REVENUE & ROAS FORECASTER", "Econometric Regression", "research/models/analytics/revenue_forecaster.pkl", "Scikit-Learn 1.8", "Scaled Campaign Features", "[ROAS, CAC, CVR]", "2.1ms", "R² = 0.894", EMERALD_GREEN),
        ("BRAND VOICE CLASSIFIER", "Text Scoring", "research/models/content/brand_voice_classifier.pkl", "Scikit-Learn 1.8", "TF-IDF + Text Statistics", "Quality Score [0-10]", "3.4ms", "MSE = 0.12", PURPLE_ACCENT),
        ("CLIP-ViT QUALITY REGRESSOR", "Zero-Shot Vision", "research/models/cv/creative_quality_regressor.pkl", "CLIP-ViT (ONNX)", "512-dim ViT Embeddings", "Aesthetic Score [0-10]", "4.8ms", "Accuracy = 91.2%", AMBER_WARNING),
        ("BGE VECTOR EMBEDDINGS", "Dense Semantic Memory", "storage/qdrant_rag / BAAI/bge-small-en-v1.5", "FastEmbed BGE", "Variable Text Chunks", "384-dim Dense Vector", "23.3ms", "MRR = 1.00", BLUE_SUPPORT),
        ("DIRICHLET CONSTRAINT GUARD", "Safety Projection", "src/adpilot/agents/optimizer/dirichlet_guard.py", "Pure Python/NumPy", "Raw Suggested Actions", "Simplex Projected Weights", "0.2ms", "0 Violations", CYAN_DARK)
    ]

    for i, (mname, mcat, mpath, mframe, minp, mout, mlat, macc, mcol) in enumerate(ml_models):
        c_idx = i % 3
        r_idx = i // 3
        c_left = 0.8 + c_idx * 3.98
        c_top = 1.75 + r_idx * 2.5
        draw_card(s11, c_left, c_top, 3.75, 2.35, mname, mcol)
        tbox = s11.shapes.add_textbox(Inches(c_left + 0.3), Inches(c_top + 0.5), Inches(3.45), Inches(1.9))
        tf = tbox.text_frame
        tf.word_wrap = True
        specs = [f"Category: {mcat}", f"Framework: {mframe}  •  Latency: {mlat}", f"Input: {minp}", f"Output: {mout}", f"Metric: {macc}"]
        for s_idx, spec in enumerate(specs):
            p = tf.add_paragraph() if s_idx > 0 else tf.paragraphs[0]
            p.text = spec
            p.font.name = FONT_MONO if s_idx in [1, 4] else FONT_BODY
            p.font.size = Pt(8.5)
            if s_idx == 4:
                p.font.bold = True
                p.font.color.rgb = mcol
            elif s_idx == 0:
                p.font.bold = True
                p.font.color.rgb = TEXT_WHITE
            else:
                p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 12: COMPUTER VISION QUALITY GATE (With CV HUD Visual)
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    add_header_footer(s12, "VISION QUALITY GATE", "Computer Vision Quality Gate & Brand Safety", "Automated multi-modal inspection using CLIP-ViT and mathematical edge detection to guarantee zero creative defects prior to publishing.", 12)
    
    # Left Column: 3 Inspection Pillars (5.8 inches)
    cv_items = [
        ("1. CLIP-ViT AESTHETIC SCORING", "Cosine similarity against aesthetic benchmark vectors. Threshold ≥ 7.0/10 required for auto-approval.", CYAN_PRIMARY),
        ("2. SAFE-ZONE MARGIN AUDIT", "Mathematical boundary mask inspection verifying that text and logos avoid mobile UI occlusion zones.", PURPLE_ACCENT),
        ("3. WCAG AAA CONTRAST RATIO", "Luminance ratio calculations guaranteeing contrast ≥ 7.0:1 for body copy and ≥ 4.5:1 for headers.", EMERALD_GREEN)
    ]
    for i, (title, desc, col) in enumerate(cv_items):
        ct = 1.75 + i * 1.62
        draw_card(s12, 0.8, ct, 5.75, 1.48, title, col)
        tbox = s12.shapes.add_textbox(Inches(1.0), Inches(ct + 0.55), Inches(5.25), Inches(0.95))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = FONT_BODY
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_MUTED

    # Right Column: CV HUD Visual (5.9 inches)
    embed_framed_image(
        s12,
        left=6.75,
        top=1.75,
        width=5.78,
        height=4.9,
        image_path=get_asset("cv_gate.jpg"),
        caption="Real-Time Computer Vision Inspection: Safe-Zones, Heatmaps & Contrast",
        border_color=CYAN_PRIMARY
    )

    # =========================================================================
    # SLIDE 13: NANO BANANA CREATIVE STUDIO (With Studio UI Graphic)
    # =========================================================================
    s13 = prs.slides.add_slide(blank_layout)
    add_header_footer(s13, "CREATIVE STUDIO", "Nano Banana Studio: Multi-Aspect Ratio Creative Synthesis", "Multi-modal visual asset generator translating structured strategy briefs into multi-format diffusion directives with extracted brand color tokens.", 13)
    
    # Left Column: 4 Formats (5.8 inches)
    formats = [
        ("16:9 LANDSCAPE BANNER", "1920x1080 • Google Display, YouTube & Desktop Banners", CYAN_PRIMARY),
        ("1:1 SQUARE FEED", "1080x1080 • Meta Carousel, Instagram Grid & LinkedIn Feed", EMERALD_GREEN),
        ("4:5 PORTRAIT FEED", "1080x1350 • Instagram Mobile Feed (Maximum screen real estate)", PURPLE_ACCENT),
        ("9:16 STORY & REELS", "1080x1920 • Instagram Stories, TikTok Ads & YouTube Shorts", AMBER_WARNING)
    ]
    for i, (title, desc, col) in enumerate(formats):
        ct = 1.75 + i * 1.22
        draw_card(s13, 0.8, ct, 5.75, 1.08, title, col)
        tbox = s13.shapes.add_textbox(Inches(1.0), Inches(ct + 0.5), Inches(5.25), Inches(0.65))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = FONT_BODY
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_MUTED

    # Right Column: Studio UI Visual (5.9 inches)
    embed_framed_image(
        s13,
        left=6.75,
        top=1.75,
        width=5.78,
        height=4.9,
        image_path=get_asset("nano_studio.jpg"),
        caption="Nano Banana Studio: Multi-Aspect Ratio Generation & Brand Swatches",
        border_color=PURPLE_ACCENT
    )

    # =========================================================================
    # SLIDE 14: DUAL-STREAM HYBRID RAG & MEMORY (With RAG Diagram)
    # =========================================================================
    s14 = prs.slides.add_slide(blank_layout)
    add_header_footer(s14, "RAG & MEMORY", "Dual-Stream Hybrid RAG & Multi-Tier Agent Memory", "Eliminating hallucinations by fusing dense semantic embeddings with sparse keyword search and 4-tier persistent memory.", 14)
    
    # Left Column: Math & 4-Tier Memory (5.8 inches)
    draw_card(s14, 0.8, 1.75, 5.75, 4.9, "Hybrid RRF Retrieval & Memory Hierarchy", CYAN_PRIMARY)
    tbox_l = s14.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.25), Inches(4.3))
    tf_l = tbox_l.text_frame
    tf_l.word_wrap = True

    rag_text = [
        ("MATHEMATICAL RRF FUSION FORMULA", "RRF(d ∈ D) = ∑_{m ∈ M} ( w_m / ( k + r_m(d) ) )", CYAN_PRIMARY),
        ("STREAM A: BAAI/bge-small-en-v1.5", "384-dimensional dense semantic vector similarity in Qdrant.", TEXT_WHITE),
        ("STREAM B: BM25 LEXICAL SEARCH", "Exact keyword matching for technical product terms and SKU codes.", TEXT_WHITE),
        ("4-TIER PERSISTENT MEMORY", "• Tier 1: In-Memory Working Session Cache (0.2ms)\n• Tier 2: Persistent Brand Identity SQLite (1.1ms)\n• Tier 3: Persona Vector Collections in Qdrant (4.2ms)\n• Tier 4: PPO Trajectory Replay Buffer", EMERALD_GREEN)
    ]
    for k, v, col in rag_text:
        p1 = tf_l.add_paragraph()
        p1.text = k
        p1.font.name = FONT_MONO
        p1.font.size = Pt(9)
        p1.font.bold = True
        p1.font.color.rgb = col
        p2 = tf_l.add_paragraph()
        p2.text = v
        p2.font.name = FONT_BODY
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = TEXT_MUTED
        tf_l.add_paragraph().text = ""

    # Right Column: RAG Architecture Diagram (5.9 inches)
    embed_framed_image(
        s14,
        left=6.75,
        top=1.75,
        width=5.78,
        height=4.9,
        image_path=get_asset("rag_hybrid_retrieval.png"),
        caption="Dual-Stream Dense + Sparse Hybrid Retrieval Architecture",
        border_color=CYAN_PRIMARY
    )

    # =========================================================================
    # SLIDE 15: REINFORCEMENT LEARNING: CONTINUOUS PPO (With 3D Landscape Visual)
    # =========================================================================
    s15 = prs.slides.add_slide(blank_layout)
    add_header_footer(s15, "REINFORCEMENT LEARNING", "Continuous PPO Actor-Critic Policy Architecture", "Autonomous budget optimization using Proximal Policy Optimization (PPO) over a 12-dimensional continuous marketing state space.", 15)
    
    # Left Column: Mathematical Formulation (5.8 inches)
    draw_card(s15, 0.8, 1.75, 5.75, 4.9, "Mathematical Policy & Objective Formulation", PURPLE_ACCENT)
    tbox_p = s15.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.25), Inches(4.3))
    tf_p = tbox_p.text_frame
    tf_p.word_wrap = True

    p_lines = [
        ("1. STATE SPACE S ∈ ℝ¹²", "Weights, ROAS, CAC, CTR, CVR, Quality, Budget, Duration, Market Penetration, Creative Diversity.", CYAN_PRIMARY),
        ("2. ACTION SPACE a_t ∈ Δ^K", "Dirichlet Simplex bounds: ∑ a_k = 1.0, floor a_k ≥ 0.05, ceiling a_k ≤ 0.65.", TEXT_WHITE),
        ("3. CLIPPED SURROGATE LOSS", "L^{CLIP}(θ) = Ê_t [ min( r_t(θ)Â_t, clip(r_t(θ), 1-ε, 1+ε)Â_t ) ]\nClipping parameter ε = 0.2 ensures policy update stability.", EMERALD_GREEN),
        ("4. REWARD FUNCTION R(s, a)", "R(s, a) = w₁ · ΔROAS - w₂ · ΔCAC - λ · RiskPenalty\nw₁ = 1.0, w₂ = 0.5, λ = 0.2.", AMBER_WARNING)
    ]
    for k, v, col in p_lines:
        p1 = tf_p.add_paragraph()
        p1.text = k
        p1.font.name = FONT_MONO
        p1.font.size = Pt(9)
        p1.font.bold = True
        p1.font.color.rgb = col
        p2 = tf_p.add_paragraph()
        p2.text = v
        p2.font.name = FONT_BODY
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = TEXT_MUTED
        tf_p.add_paragraph().text = ""

    # Right Column: 3D Policy Gradient Landscape Visual (5.9 inches)
    embed_framed_image(
        s15,
        left=6.75,
        top=1.75,
        width=5.78,
        height=4.9,
        image_path=get_asset("rl_ppo.jpg"),
        caption="Actor-Critic Policy Gradient Landscape & State Trajectory Convergence",
        border_color=PURPLE_ACCENT
    )

    # =========================================================================
    # SLIDE 16: RL BUSINESS CASE: DYNAMIC BUDGET REALLOCATION (With Training Curve)
    # =========================================================================
    s16 = prs.slides.add_slide(blank_layout)
    add_header_footer(s16, "PPO CASE STUDY", "RL Optimization Simulation: Dynamic Budget Rebalancing", "Demonstrating PPO Actor-Critic policy action selection in a controlled econometric marketing simulation environment.", 16)
    
    # Left Column: Before vs Action vs Outcome (5.8 inches)
    b_cases = [
        ("1. BASELINE (STATIC ALLOCATION)", "LinkedIn: 35% ($3.5k) • Meta: 25% ($2.5k) • Google: 40% ($4.0k)\nPredicted ROAS: 3.20x  •  CAC: $49.00  •  Waste Ratio: 24%", ROSE_DANGER),
        ("2. PPO REBALANCE ACTION a_t", "LinkedIn: +12% ($5,700) • Meta: Hold ($3,500) • Google: -12% ($800)\nDirichlet guard preserves simplex (Sum = 100%, 0 violations).", CYAN_PRIMARY),
        ("3. SIMULATED OUTCOME", "Blended ROAS: 3.84x (+20.0%)  •  Blended CAC: $42.10 (-14.1%)\nExpected Reward R(s, a) = +0.48 Net Gain  [CONTROLLED SIMULATION]", EMERALD_GREEN)
    ]
    for i, (title, desc, col) in enumerate(b_cases):
        ct = 1.75 + i * 1.62
        draw_card(s16, 0.8, ct, 5.75, 1.48, title, col)
        tbox = s16.shapes.add_textbox(Inches(1.0), Inches(ct + 0.5), Inches(5.25), Inches(0.98))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = FONT_BODY
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MUTED

    # Right Column: RL Training Convergence Graph (5.9 inches)
    embed_framed_image(
        s16,
        left=6.75,
        top=1.75,
        width=5.78,
        height=4.9,
        image_path=get_asset("rl_optimizer_loss.png"),
        caption="PPO Policy Actor-Critic Loss & Reward Convergence Curve",
        border_color=EMERALD_GREEN
    )

    # =========================================================================
    # SLIDE 17: HUMAN-IN-THE-LOOP (HITL) GOVERNANCE
    # =========================================================================
    s17 = prs.slides.add_slide(blank_layout)
    add_header_footer(s17, "GOVERNANCE & TRUST", "Human-in-the-Loop (HITL) Governance & Cryptography", "Cryptographic review gates enforcing strict enterprise oversight for high-risk autonomous decisions, budget rebalancing, and live ad publishing.", 17)
    
    gov_cards = [
        ("1. RISK-BASED DECISION GATES", "Automated Triage by Impact", [("Low Risk (Auto-Execute)", "Headline A/B variations, minor keyword adjustments, standard budget allocations."), ("Medium Risk (Flagged)", "Budget reallocations > 10%, new channel activation, audience expansion."), ("High Risk (Gate Lock)", "Live ad spend dispatch, brand repositioning, total budget redistribution.")], CYAN_PRIMARY),
        ("2. ROLE-BASED ACCESS (RBAC)", "Multi-Stakeholder Authority", [("Campaign Director", "Full authority to approve multi-channel strategy, budget shifts, and live dispatch."), ("Brand Auditor", "Audits CLIP-ViT visual quality scores, tone compliance, and copyright safety."), ("Growth Lead", "Monitors PPO reward trajectories, CAC thresholds, and channel attribution.")], PURPLE_ACCENT),
        ("3. HMAC-SHA256 AUDIT LOGS", "Immutable Cryptographic Ledger", [("Formula", "HMAC(K, m) = H( (K' ⊕ opad) || H( (K' ⊕ ipad) || m ) )"), ("Non-Repudiation", "Every human approval generates a cryptographic SHA-256 signature token."), ("Audit Trail", "Append-only database ledger guarantees complete forensic traceability.")], AMBER_WARNING)
    ]

    for i, (title, subtitle, items, col) in enumerate(gov_cards):
        c_left = 0.8 + i * 3.98
        draw_card(s17, c_left, 1.75, 3.75, 4.9, title, col)
        tbox = s17.shapes.add_textbox(Inches(c_left + 0.3), Inches(2.2), Inches(3.25), Inches(4.3))
        tf = tbox.text_frame
        tf.word_wrap = True
        p_sub = tf.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.name = FONT_HEADING
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_WHITE
        tf.add_paragraph().text = ""
        for k, v in items:
            p_k = tf.add_paragraph()
            p_k.text = k.upper()
            p_k.font.name = FONT_MONO
            p_k.font.size = Pt(8.5)
            p_k.font.bold = True
            p_k.font.color.rgb = col
            p_v = tf.add_paragraph()
            p_v.text = v
            p_v.font.name = FONT_BODY
            p_v.font.size = Pt(8.5)
            p_v.font.color.rgb = TEXT_MUTED
            tf.add_paragraph().text = ""

    # =========================================================================
    # SLIDE 18: END-TO-END CAMPAIGN SIMULATION WORKFLOW
    # =========================================================================
    s18 = prs.slides.add_slide(blank_layout)
    add_header_footer(s18, "LIVE DEMO FLOW", "End-to-End Campaign Execution Simulation Trace", "Complete walkthrough of how a single structured campaign brief flows through all autonomous intelligence, creative, and governance stages.", 18)
    
    sim_steps = [
        ("Step 01", "Brief Ingestion", "Client enters 'VisionGuard AI' ($10k, 30 days, B2B Security). Pydantic validates contract.", CYAN_PRIMARY),
        ("Step 02", "Strategy Synthesis", "GPT-4o Router formulates 3-stage funnel: LinkedIn (45%), Meta (35%), Google (20%).", BLUE_SUPPORT),
        ("Step 03", "Hybrid Vector RAG", "FastEmbed BGE retrieves ICP pain points ('Alert fatigue, Compliance audit'). MRR: 1.0.", PURPLE_ACCENT),
        ("Step 04", "Content Copywriting", "Claude 3.5 Sonnet drafts 8 ad variants + 4-stage email nurture. Ridge score: 5.43.", CYAN_DARK),
        ("Step 05", "Nano Banana Studio", "Synthesizes multi-aspect prompts (16:9, 1:1, 4:5, 9:16) with obsidian & cyan palette.", EMERALD_GREEN),
        ("Step 06", "CLIP-ViT Vision Gate", "Zero-shot visual inspection passes: 8.7/10 quality, safe zone margin & contrast compliant.", AMBER_WARNING),
        ("Step 07", "Predictive Analytics", "Ridge regressor forecasts ROAS: 3.84x, CAC: $42.10, Health Score: 87.5/100.", ROSE_DANGER),
        ("Step 08", "PPO Policy Rebalance", "Continuous policy reallocates budget (+12% LinkedIn) subject to Dirichlet simplex guards.", CYAN_PRIMARY),
        ("Step 09", "HITL Sign-Off", "Campaign Director approves rebalance. Cryptographic SHA-256 signature committed to ledger.", EMERALD_GREEN),
        ("Step 10", "Safe Dispatch", "Idempotent multi-network dry-run dispatch executed. WebSocket streams telemetry to UI.", PURPLE_ACCENT)
    ]

    for i, (step_num, step_name, step_desc, col) in enumerate(sim_steps):
        col_idx = i % 5
        row_idx = i // 5
        c_left = 0.8 + col_idx * 2.38
        c_top = 1.75 + row_idx * 2.5
        draw_card(s18, c_left, c_top, 2.22, 2.35, f"{step_num}: {step_name}", col)
        tbox = s18.shapes.add_textbox(Inches(c_left + 0.3), Inches(c_top + 0.4), Inches(1.98), Inches(1.85))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step_desc
        p.font.name = FONT_BODY
        p.font.size = Pt(8.5)
        p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 19: EMPIRICAL SIMULATION RESULTS & BENCHMARKS
    # =========================================================================
    s19 = prs.slides.add_slide(blank_layout)
    add_header_footer(s19, "EMPIRICAL BENCHMARKS", "Empirical Evaluation & Comparative Benchmark Analysis", "Rigorous comparison between traditional manual agencies, monolithic single LLM prompt wrappers, and the ADPilot Pro Multi-Agent OS.", 19)
    
    comp_headers = ["EVALUATION METRIC", "TRADITIONAL AGENCY", "SINGLE LLM CHATBOT", "ADPILOT PRO MULTI-AGENT OS", "ADVANTAGE"]
    comp_rows = [
        ["Brief-to-Launch Latency", "10–14 Days", "2–4 Hours", "< 15 Seconds (Compute)", "100x Faster Launch"],
        ["Context Integrity Loss", "40% Handoff Decay", "28% Context Drift", "0% (Immutable Pydantic v2)", "Zero Drift Guard"],
        ["Simulated ROAS Output", "3.20x Baseline", "3.35x Unbounded", "3.84x (PPO Optimized)", "+20.0% Efficiency"],
        ["Creative Formats Supported", "Manual Resizing", "Text Prompts Only", "Multi-Aspect (16:9, 1:1, 4:5, 9:16)", "Full Omnichannel"],
        ["Visual Safety Inspection", "Manual Review", "None", "CLIP-ViT (100% Margin Check)", "Automated Safety"],
        ["Closed-Loop Budget Tuning", "Weekly Human Edits", "None", "Continuous PPO Actor-Critic", "Real-Time Tuning"],
        ["Governance Auditability", "Scattered Emails", "None", "Cryptographic HMAC-SHA256", "Enterprise Security"]
    ]

    t_shape = s19.shapes.add_table(8, 5, Inches(0.8), Inches(1.75), Inches(11.733), Inches(4.9))
    tbl = t_shape.table
    tbl.columns[0].width = Inches(2.3)
    tbl.columns[1].width = Inches(2.1)
    tbl.columns[2].width = Inches(2.1)
    tbl.columns[3].width = Inches(3.2)
    tbl.columns[4].width = Inches(2.033)

    for r_idx, row in enumerate([comp_headers] + comp_rows):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.fill.solid()
            if r_idx == 0:
                cell.fill.fore_color.rgb = RGBColor(16, 24, 45)
            else:
                cell.fill.fore_color.rgb = BG_CARD if r_idx % 2 == 0 else BG_CARD_SEC

            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = val
            p.font.name = FONT_MONO if (r_idx == 0 or c_idx in [0, 4]) else FONT_BODY
            p.font.size = Pt(8.5) if r_idx > 0 else Pt(9)
            if r_idx == 0:
                p.font.bold = True
                p.font.color.rgb = CYAN_PRIMARY
            else:
                if c_idx == 3:
                    p.font.bold = True
                    p.font.color.rgb = CYAN_PRIMARY
                elif c_idx == 4:
                    p.font.bold = True
                    p.font.color.rgb = EMERALD_GREEN
                elif c_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = TEXT_WHITE
                else:
                    p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 20: PRODUCTION DASHBOARD & UI EXPERIENCE (With Dashboard Preview)
    # =========================================================================
    s20 = prs.slides.add_slide(blank_layout)
    add_header_footer(s20, "PRODUCT EXPERIENCE", "Production Web Application & Cyber UI Dashboard", "ADPilot Pro is a live, working enterprise software application featuring a React 18 single-page application and FastAPI backend.", 20)
    
    # Left Column: Key Features (5.8 inches)
    ui_feats = [
        ("EXECUTIVE KPI VIEW", "Live spend, blended AI ROAS (3.84x), CAC ($42.10), and glowing SVG ROAS curves.", CYAN_PRIMARY),
        ("INTERACTIVE 18-STAGE DAG", "Node graph with live beacons, model tags, latency meters, and schema inspection.", BLUE_SUPPORT),
        ("NANO BANANA STUDIO", "Multi-aspect gallery (16:9, 1:1, 4:5, 9:16) with WCAG AAA badges & hex copiers.", PURPLE_ACCENT),
        ("HITL GOVERNANCE CENTER", "RBAC role switcher with HMAC-SHA256 append-only cryptographic audit ledger.", AMBER_WARNING)
    ]
    for i, (title, desc, col) in enumerate(ui_feats):
        ct = 1.75 + i * 1.22
        draw_card(s20, 0.8, ct, 5.75, 1.08, title, col)
        tbox = s20.shapes.add_textbox(Inches(1.0), Inches(ct + 0.5), Inches(5.25), Inches(0.65))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = FONT_BODY
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_MUTED

    # Right Column: Dashboard UI Preview (5.9 inches)
    embed_framed_image(
        s20,
        left=6.75,
        top=1.75,
        width=5.78,
        height=4.9,
        image_path=get_asset("001.png") or get_asset("dashboard_preview.png") or get_asset("ui_predesign.png"),
        caption="Live React 18 Enterprise Marketing Cyber Dashboard",
        border_color=CYAN_PRIMARY
    )

    # =========================================================================
    # SLIDE 21: FULL-STACK SYSTEM ARCHITECTURE
    # =========================================================================
    s21 = prs.slides.add_slide(blank_layout)
    add_header_footer(s21, "SYSTEM ARCHITECTURE", "Full-Stack Enterprise Engineering Architecture", "Layered microservice architecture designed for high-concurrency asynchronous execution, data isolation, and low-latency inference.", 21)
    
    arch_layers = [
        ("1. PRESENTATION TIER", "React 18 SPA  •  Vite  •  Tailwind CSS  •  Zustand State  •  Lucide Icons", CYAN_PRIMARY),
        ("2. API & GATEWAY TIER", "FastAPI (Async Python 3.12)  •  OpenAPI / Swagger  •  JWT Auth  •  Pydantic v2", BLUE_SUPPORT),
        ("3. ORCHESTRATION TIER", "Master 18-Stage Execution DAG  •  AsyncIO Gather  •  AgentEventBus  •  WebSocket Telemetry", PURPLE_ACCENT),
        ("4. MULTI-AGENT FLEET", "18 Specialized Agents extending BaseAgent  •  Pydantic Immutable Contracts", EMERALD_GREEN),
        ("5. NEURAL & ML ENGINES", "PyTorch PPO Policy (.pt)  •  Scikit-Learn Ridge (.pkl)  •  CLIP-ViT (ONNX)", AMBER_WARNING),
        ("6. RAG & MEMORY STORE", "FastEmbed BGE (384-d)  •  Qdrant Vector DB  •  4-Tier Multi-Tier Memory Store", CYAN_DARK),
        ("7. PERSISTENCE & WORKERS", "Async SQLAlchemy  •  SQLite / PostgreSQL  •  ARQ In-Process Task Workers", ROSE_DANGER),
        ("8. GOVERNANCE & AUDIT", "HMAC-SHA256 Append-Only Audit Ledger  •  Dirichlet Simplex Bounds", CYAN_PRIMARY)
    ]

    for i, (layer_name, tech_stack, col) in enumerate(arch_layers):
        l_top = 1.75 + i * 0.62
        draw_card(s21, 0.8, l_top, 11.733, 0.55, border_color=RGBColor(30, 41, 59))
        tbox = s21.shapes.add_textbox(Inches(1.0), Inches(l_top + 0.1), Inches(3.5), Inches(0.35))
        tf = tbox.text_frame
        p = tf.paragraphs[0]
        p.text = layer_name
        p.font.name = FONT_MONO
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = col

        tbox_s = s21.shapes.add_textbox(Inches(4.5), Inches(l_top + 0.1), Inches(7.8), Inches(0.35))
        tf_s = tbox_s.text_frame
        p_s = tf_s.paragraphs[0]
        p_s.text = tech_stack
        p_s.font.name = FONT_BODY
        p_s.font.size = Pt(9)
        p_s.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 22: QUANTIFIABLE BUSINESS VALUE & ROI
    # =========================================================================
    s22 = prs.slides.add_slide(blank_layout)
    add_header_footer(s22, "COMMERCIAL ROI", "Quantifiable Business Value & Enterprise ROI", "Transforming marketing from a high-friction cost center into a predictable, capital-efficient, and continuously improving growth engine.", 22)
    
    roi_cards = [
        ("10x VELOCITY", "Launch Time Compression", "Compresses the full brief-to-launch cycle from 10–14 days to under 15 seconds of compute time.", CYAN_PRIMARY),
        ("60%+ COST DROP", "Overhead Elimination", "Dramatically slashes agency fees and manual copy-pasting hours spent on routine formatting.", EMERALD_GREEN),
        ("24/7 OPTIMIZATION", "Continuous RL Bidding", "PPO Actor-Critic policy continuously tunes channel budgets to capture high-performing ad windows.", PURPLE_ACCENT),
        ("100% BRAND GUARD", "Zero Context Drift", "Immutable brief propagation guarantees complete brand tone consistency across all marketing assets.", AMBER_WARNING)
    ]

    for i, (val, label, desc, col) in enumerate(roi_cards):
        c_left = 0.8 + i * 2.98
        draw_card(s22, c_left, 1.75, 2.78, 4.9, val, col)
        tbox = s22.shapes.add_textbox(Inches(c_left + 0.3), Inches(2.2), Inches(2.28), Inches(4.3))
        tf = tbox.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = label
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE
        tf.add_paragraph().text = ""
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 23: TECHNICAL RIGOR & VERIFIED TESTING
    # =========================================================================
    s23 = prs.slides.add_slide(blank_layout)
    add_header_footer(s23, "EVALUATION & RIGOR", "Technical Rigor, Testing & Quantitative Verification", "Every component of ADPilot Pro has been verified through automated test suites, contract schema audits, and end-to-end integration tests.", 23)
    
    eval_metrics = [
        ("52 / 52", "FRONTEND VITEST TESTS", "100% passing across Component UI, Brief Form, Design Preview, Result Display.", EMERALD_GREEN),
        ("269 / 269", "BACKEND PYTEST TESTS", "100% passing across 18 Agents, BaseAgent contract inheritance, Pydantic v2 schemas.", CYAN_PRIMARY),
        ("0 ERRORS", "TYPESCRIPT COMPILATION", "Clean `tsc --noEmit` build across all React 18 interfaces and typed state hooks.", PURPLE_ACCENT),
        ("< 0.1%", "SCHEMA VALIDATION ERROR", "SVER < 0.1% across multi-agent handoffs, confirming zero payload corruption.", AMBER_WARNING)
    ]

    for i, (val, label, desc, col) in enumerate(eval_metrics):
        c_left = 0.8 + i * 2.98
        draw_card(s23, c_left, 1.75, 2.78, 2.35, border_color=col)
        vbox = s23.shapes.add_textbox(Inches(c_left + 0.3), Inches(1.95), Inches(2.28), Inches(0.75))
        tf = vbox.text_frame
        p = tf.paragraphs[0]
        p.text = val
        p.font.name = FONT_HEADING
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = col

        lbox = s23.shapes.add_textbox(Inches(c_left + 0.3), Inches(2.7), Inches(2.28), Inches(1.3))
        tf_l = lbox.text_frame
        tf_l.word_wrap = True
        p1 = tf_l.paragraphs[0]
        p1.text = label
        p1.font.name = FONT_MONO
        p1.font.size = Pt(9)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE
        p2 = tf_l.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = TEXT_MUTED

    draw_card(s23, 0.8, 4.35, 11.733, 2.3, "Automated Verification Framework", CYAN_PRIMARY)
    tbox_t = s23.shapes.add_textbox(Inches(1.0), Inches(4.75), Inches(11.333), Inches(1.8))
    tf_t = tbox_t.text_frame
    tf_t.word_wrap = True
    t_bullets = [
        "• Unit & Contract Tests: Verifies immutable data schemas, Pydantic type coercion, and tool mocking for offline execution.",
        "• Integration DAG Verification: 18 verification scripts (verify_phase1.py through verify_phase16.py) validating pipeline execution.",
        "• Mathematical Convergence: PPO Actor-Critic policy gradient updates tested against simulated baseline environments.",
        "• UI End-to-End Build: Clean production build in 9.78s transforming 1,503 Vite modules with zero bundle warnings."
    ]
    for bullet in t_bullets:
        p = tf_t.add_paragraph()
        p.text = bullet
        p.font.name = FONT_BODY
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 24: VERTICAL ICP FIT & INDUSTRY APPLICATIONS
    # =========================================================================
    s24 = prs.slides.add_slide(blank_layout)
    add_header_footer(s24, "MARKET ADOPTION", "Multi-Vertical Market Application & ICP Fit", "The Product Classifier Agent dynamically adapts funnel frameworks, tone profiles, and creative tokens across major commercial verticals.", 24)
    
    verticals = [
        ("B2B SAAS & TECH", "High-LTV Customer Acquisition", [("Funnel Focus", "Multi-stage lead nurture & demo requests"), ("Channels", "LinkedIn Sponsored, Google Search RSAs"), ("Agent Action", "Deep technical battlecards & ROI case studies")], CYAN_PRIMARY),
        ("E-COMMERCE & DTC", "High-Volume Transactional Scale", [("Funnel Focus", "Direct checkout conversion & ROAS scale"), ("Channels", "Meta Carousel, Instagram Feed, TikTok"), ("Agent Action", "Nano Banana multi-variant creative diffusion")], EMERALD_GREEN),
        ("DIGITAL AGENCIES", "Multi-Client Operational Scale", [("Funnel Focus", "10x throughput per marketing strategist"), ("Channels", "Omnichannel multi-client client workspaces"), ("Agent Action", "Automated white-label strategy playbook generation")], PURPLE_ACCENT),
        ("FINTECH & HEALTH", "High-Trust Regulated Sectors", [("Funnel Focus", "Trust building, compliance & qualification"), ("Channels", "High-intent search, compliance landing pages"), ("Agent Action", "Strict HITL review & cryptographic audit trails")], AMBER_WARNING)
    ]

    for i, (title, subtitle, kvs, col) in enumerate(verticals):
        c_left = 0.8 + i * 2.98
        draw_card(s24, c_left, 1.75, 2.78, 4.9, title, col)
        tbox = s24.shapes.add_textbox(Inches(c_left + 0.3), Inches(2.2), Inches(2.28), Inches(4.3))
        tf = tbox.text_frame
        tf.word_wrap = True
        p_sub = tf.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.name = FONT_HEADING
        p_sub.font.size = Pt(10)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_WHITE
        tf.add_paragraph().text = ""
        for k, v in kvs:
            p_k = tf.add_paragraph()
            p_k.text = k.upper()
            p_k.font.name = FONT_MONO
            p_k.font.size = Pt(8.5)
            p_k.font.bold = True
            p_k.font.color.rgb = col
            p_v = tf.add_paragraph()
            p_v.text = v
            p_v.font.name = FONT_BODY
            p_v.font.size = Pt(8.5)
            p_v.font.color.rgb = TEXT_MUTED
            tf.add_paragraph().text = ""

    # =========================================================================
    # SLIDE 25: COMMERCIAL SAAS REVENUE MODEL
    # =========================================================================
    s25 = prs.slides.add_slide(blank_layout)
    add_header_footer(s25, "BUSINESS MODEL", "Commercial SaaS Monetization & Revenue Architecture", "Scalable subscription tiers paired with usage-based neural compute credits for high gross-margin enterprise monetization.", 25)
    
    tiers = [
        ("GROWTH SAAS", "$499 / mo", "SMBs & DTC Brands", ["• Full 18-Agent Autonomous Pipeline", "• Up to 20 Campaigns / Month", "• Nano Banana Visual Asset Studio", "• CLIP-ViT Pre-Flight Quality Check", "• Standard Email & Chat Support"], CYAN_PRIMARY),
        ("AGENCY PRO", "$1,499 / mo", "Digital Marketing Agencies", ["• Multi-Client Isolated Workspaces", "• Unlimited Autonomous Campaigns", "• Custom Brand Tone & Hex Profiles", "• White-Label PDF Playbook Export", "• Dedicated Technical Account Manager"], PURPLE_ACCENT),
        ("ENTERPRISE CUSTOM", "$5,000+ / mo", "Large Enterprises & Scaleups", ["• Dedicated VPC On-Prem Deployment", "• Custom Ad Network Bidding API Sync", "• Fine-Tuned Custom PPO Policy Models", "• Private Isolated Qdrant Vector Memory", "• 99.9% Uptime SLA & 24/7 Security Support"], EMERALD_GREEN)
    ]

    for i, (tier_name, price, target, features_list, col) in enumerate(tiers):
        c_left = 0.8 + i * 3.98
        draw_card(s25, c_left, 1.75, 3.75, 4.9, tier_name, col)
        tbox = s25.shapes.add_textbox(Inches(c_left + 0.3), Inches(2.2), Inches(3.25), Inches(4.3))
        tf = tbox.text_frame
        tf.word_wrap = True
        p_pr = tf.paragraphs[0]
        p_pr.text = price
        p_pr.font.name = FONT_HEADING
        p_pr.font.size = Pt(22)
        p_pr.font.bold = True
        p_pr.font.color.rgb = col
        p_tg = tf.add_paragraph()
        p_tg.text = target
        p_tg.font.name = FONT_MONO
        p_tg.font.size = Pt(9.5)
        p_tg.font.bold = True
        p_tg.font.color.rgb = TEXT_WHITE
        tf.add_paragraph().text = ""
        for feat in features_list:
            p = tf.add_paragraph()
            p.text = feat
            p.font.name = FONT_BODY
            p.font.size = Pt(9.5)
            p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 26: SYSTEM LIMITATIONS & SCIENTIFIC TRANSPARENCY
    # =========================================================================
    s26 = prs.slides.add_slide(blank_layout)
    add_header_footer(s26, "SCIENTIFIC RIGOR", "System Limitations & Technical Transparency", "Maintaining rigorous academic honesty by explicitly separating current verified capabilities from current environmental constraints.", 26)
    
    limitations = [
        ("1. SIMULATION ENVIRONMENT", "PPO Policy Simulator Bounds", [("Current Status", "Trained on econometric market simulators with realistic CTR/CPC noise distributions."), ("Limitation", "Real-world market dynamics feature stochastic external shocks not fully modeled in synthetic replay buffers."), ("Mitigation", "Dirichlet simplex bounding constraints prevent catastrophic budget shifts in unobserved states.")], AMBER_WARNING),
        ("2. VIDEO SYNTHESIS PIPELINE", "Multi-Modal Creative Scope", [("Current Status", "Generates high-resolution multi-aspect static image banners and rich copywriting."), ("Limitation", "Direct autonomous video rendering (e.g. Sora / Runway API) is currently in architectural design phase."), ("Mitigation", "Generates detailed storyboard prompt scripts ready for external video production pipelines.")], PURPLE_ACCENT),
        ("3. API RATE LIMITS & TOKENS", "External LLM Dependencies", [("Current Status", "Multi-provider routing between OpenAI GPT-4o and Anthropic Claude 3.5 Sonnet."), ("Limitation", "High concurrency runs can trigger third-party API rate limiting and token budget inflation."), ("Mitigation", "Implemented FastEmbed BGE vector caching and local mock simulation fallback for offline execution.")], CYAN_PRIMARY)
    ]

    for i, (title, subtitle, items, col) in enumerate(limitations):
        c_left = 0.8 + i * 3.98
        draw_card(s26, c_left, 1.75, 3.75, 4.9, title, col)
        tbox = s26.shapes.add_textbox(Inches(c_left + 0.3), Inches(2.2), Inches(3.25), Inches(4.3))
        tf = tbox.text_frame
        tf.word_wrap = True
        p_sub = tf.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.name = FONT_HEADING
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_WHITE
        tf.add_paragraph().text = ""
        for k, v in items:
            p_k = tf.add_paragraph()
            p_k.text = k.upper()
            p_k.font.name = FONT_MONO
            p_k.font.size = Pt(8.5)
            p_k.font.bold = True
            p_k.font.color.rgb = col
            p_v = tf.add_paragraph()
            p_v.text = v
            p_v.font.name = FONT_BODY
            p_v.font.size = Pt(8.5)
            p_v.font.color.rgb = TEXT_MUTED
            tf.add_paragraph().text = ""

    # =========================================================================
    # SLIDE 27: STRATEGIC FUTURE ROADMAP
    # =========================================================================
    s27 = prs.slides.add_slide(blank_layout)
    add_header_footer(s27, "PRODUCT ROADMAP", "Strategic Product Roadmap & Scale Horizons", "Clear separation between verified implemented features, active platform components, and future autonomous scaling horizons.", 27)
    
    phases = [
        ("PHASE 1: FOUNDATION", "[VERIFIED & IMPLEMENTED]", ["• 18-stage deterministic DAG execution engine", "• BaseAgent contract pattern & Pydantic v2 validation", "• GPT-4o & Claude 3.5 multi-provider intelligence routing", "• FastEmbed BGE + Qdrant Dual-Stream Hybrid RAG", "• PyTorch PPO Actor-Critic budget policy network", "• CLIP-ViT zero-shot visual aesthetic scoring", "• 52/52 Vitest & 269 Pytest automated test coverage"], EMERALD_GREEN),
        ("PHASE 2: PLATFORM", "[ACTIVE / PROTOTYPE]", ["• FastAPI asynchronous REST gateway & WebSockets", "• React 18 single-page cyber dashboard UI", "• Nano Banana multi-aspect ratio creative studio", "• SQLite & In-Memory 4-tier persistent memory", "• HMAC-SHA256 cryptographic audit trail ledger", "• Multi-role HITL review governance center", "• Dry-run multi-channel publishing adapters"], CYAN_PRIMARY),
        ("PHASE 3: SCALE & RL OPS", "[FUTURE ROADMAP]", ["• Live Meta Ads & Google Ads direct API bidding sync", "• Real-time continuous online PPO policy adaptation", "• Generative multi-modal AI video ad diffusion", "• Federated multi-tenant cross-organization learning", "• Automated A/B creative multivariate experimentation", "• Enterprise SSO & multi-region VPC clustering"], PURPLE_ACCENT)
    ]

    for i, (title, status, bullets, col) in enumerate(phases):
        c_left = 0.8 + i * 3.98
        draw_card(s27, c_left, 1.75, 3.75, 4.9, title, col)
        tbox = s27.shapes.add_textbox(Inches(c_left + 0.3), Inches(2.2), Inches(3.25), Inches(4.3))
        tf = tbox.text_frame
        tf.word_wrap = True
        p_st = tf.paragraphs[0]
        p_st.text = status
        p_st.font.name = FONT_MONO
        p_st.font.size = Pt(9.5)
        p_st.font.bold = True
        p_st.font.color.rgb = col
        tf.add_paragraph().text = ""
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.name = FONT_BODY
            p.font.size = Pt(9.5)
            p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 28: CAPSTONE DEFENSE CONCLUSION (With Hero Closing Visual)
    # =========================================================================
    s28 = prs.slides.add_slide(blank_layout)
    bg28 = s28.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg28.fill.solid()
    bg28.fill.fore_color.rgb = BG_DARK
    bg28.line.fill.background()

    top_bar28 = s28.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.04))
    top_bar28.fill.solid()
    top_bar28.fill.fore_color.rgb = CYAN_PRIMARY
    top_bar28.line.fill.background()

    # Left Column: Summary & Gratitude (6.2 inches)
    pill28 = s28.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.9), Inches(4.6), Inches(0.38))
    pill28.fill.solid()
    pill28.fill.fore_color.rgb = BG_CARD
    pill28.line.color.rgb = EMERALD_GREEN
    pill28.text_frame.paragraphs[0].text = "ACADEMIC CAPSTONE DEFENSE  //  JUNE 2026"
    pill28.text_frame.paragraphs[0].font.name = FONT_MONO
    pill28.text_frame.paragraphs[0].font.size = Pt(9.5)
    pill28.text_frame.paragraphs[0].font.bold = True
    pill28.text_frame.paragraphs[0].font.color.rgb = EMERALD_GREEN

    h28_box = s28.shapes.add_textbox(Inches(0.8), Inches(1.45), Inches(5.8), Inches(1.0))
    tf28 = h28_box.text_frame
    p28 = tf28.paragraphs[0]
    p28.text = "ADPILOT PRO"
    p28.font.name = FONT_HEADING
    p28.font.size = Pt(40)
    p28.font.bold = True
    p28.font.color.rgb = TEXT_WHITE

    sub28_box = s28.shapes.add_textbox(Inches(0.8), Inches(2.45), Inches(5.8), Inches(0.65))
    tf28_sub = sub28_box.text_frame
    tf28_sub.word_wrap = True
    p28_sub = tf28_sub.paragraphs[0]
    p28_sub.text = "One Strategic Brief  •  18 Autonomous Agents  •  Continuous Intelligence"
    p28_sub.font.name = FONT_BODY
    p28_sub.font.size = Pt(12)
    p28_sub.font.color.rgb = TEXT_MUTED

    summary_pillars = [
        ("SCIENTIFIC RIGOR", "Pydantic v2 typed contracts, continuous PPO policy gradients, Dual-Stream Hybrid RAG (MRR: 1.0), and CLIP-ViT visual quality gates.", CYAN_PRIMARY),
        ("ENGINEERING REALITY", "Production software platform with React 18 SPA, FastAPI backend, 52/52 Vitest tests, and 269 passing Pytest test suites.", EMERALD_GREEN),
        ("COMMERCIAL VALUE", "Compresses 14-day launch cycles to under 15 seconds with +20% simulated ROAS gains and HMAC-SHA256 security.", PURPLE_ACCENT)
    ]
    for i, (title, desc, col) in enumerate(summary_pillars):
        ct = 3.25 + i * 0.92
        draw_card(s28, 0.8, ct, 5.75, 0.82, title, col)
        tbox = s28.shapes.add_textbox(Inches(1.0), Inches(ct + 0.32), Inches(5.25), Inches(0.48))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = FONT_BODY
        p.font.size = Pt(8.5)
        p.font.color.rgb = TEXT_MUTED

    # Right Column: Closing Visual Hero (5.68 inches)
    embed_framed_image(
        s28,
        left=6.85,
        top=0.9,
        width=5.68,
        height=4.95,
        image_path=get_asset("hero_cover.jpg") or get_asset("hero_banner.png"),
        caption="From Strategic Intent to Autonomous Action — Thank You",
        border_color=EMERALD_GREEN
    )

    # Footer Gratitude
    cred_box = s28.shapes.add_textbox(Inches(0.8), Inches(6.1), Inches(11.733), Inches(0.8))
    tf_cred = cred_box.text_frame
    p_c1 = tf_cred.paragraphs[0]
    p_c1.text = "We sincerely thank the Digital Pioneers Initiative (DiGiLiANS), MCIT Egypt, and the Military Technical College (MTC) for their dedicated support."
    p_c1.font.name = FONT_BODY
    p_c1.font.size = Pt(9.5)
    p_c1.font.bold = True
    p_c1.font.color.rgb = TEXT_WHITE

    p_c2 = tf_cred.add_paragraph()
    p_c2.text = "ADPILOT PRO  •  AI & DATA SCIENCE TRACK  •  CAPSTONE DEFENSE 2026  •  QUESTIONS & DISCUSSION"
    p_c2.font.name = FONT_MONO
    p_c2.font.size = Pt(9.5)
    p_c2.font.bold = True
    p_c2.font.color.rgb = CYAN_PRIMARY

    # Save presentations with multiple fallbacks if locked by PowerPoint
    candidate_paths = [
        r"Presentation\FINAL_ADPILOT_ULTIMATE_EDITION.pptx",
        r"Presentation\FINAL_ADPILOT_CYBER_EDITION.pptx",
        r"Presentation\FINAL_ADPILOT_MASTER_DECK.pptx",
        r"Presentation\FINAL_ADPILOT_PRESENTATION_PRO.pptx",
        r"Presentation\FINAL_ADPILOT_PRESENTATION.pptx"
    ]
    saved_paths = []
    for path in candidate_paths:
        try:
            prs.save(path)
            saved_paths.append(path)
            print(f"[SUCCESS] High-Visuals Presentation generated: {path}")
        except PermissionError:
            print(f"[NOTICE] {path} is currently locked by PowerPoint. Saved to alternative path.")
        except Exception as e:
            print(f"[ERROR] Could not save to {path}: {e}")

if __name__ == "__main__":
    create_presentation()

